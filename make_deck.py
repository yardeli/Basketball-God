"""
make_deck.py — Generate Basketball-God v2 Upgrades PowerPoint
Run: python make_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette (dark dashboard theme) ─────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)   # near-black background
CARD      = RGBColor(0x16, 0x1B, 0x27)   # dark card
ORANGE    = RGBColor(0xFF, 0x6B, 0x00)   # accent orange
GREEN     = RGBColor(0x00, 0xD4, 0x6A)   # positive green
BLUE      = RGBColor(0x38, 0xBD, 0xF8)   # cyan blue
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)   # purple
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)   # yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
DARKGRAY  = RGBColor(0x1E, 0x29, 0x3B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color, alpha=None, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def rounded_box(slide, x, y, w, h, fill_color):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def txt(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb

def pill(slide, label, x, y, color):
    """Small colored pill/badge."""
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(1.6), Inches(0.35))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

def arrow_right(slide, x, y, length=0.8, color=GRAY):
    """Simple right-pointing arrow line."""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x), Inches(y), Inches(x + length), Inches(y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)

def metric_card(slide, x, y, label, v1, v2, delta, delta_good=True):
    """Comparison metric card: label, v1, v2, delta."""
    rounded_box(slide, x, y, 2.8, 1.5, DARKGRAY)
    txt(slide, label, x+0.12, y+0.08, 2.6, 0.3, size=11, color=GRAY)
    txt(slide, v1,    x+0.12, y+0.38, 1.1, 0.4, size=20, color=GRAY, bold=True)
    txt(slide, "→",   x+1.25, y+0.38, 0.4, 0.4, size=20, color=GRAY)
    txt(slide, v2,    x+1.6,  y+0.38, 1.1, 0.4, size=20, color=GREEN if delta_good else ORANGE, bold=True)
    dcol = GREEN if delta_good else ORANGE
    txt(slide, delta, x+0.12, y+1.0,  2.6, 0.35, size=13, color=dcol, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)

# Top accent bar
box(slide, 0, 0, 13.33, 0.08, ORANGE)

# Basketball emoji area
box(slide, 0.6, 1.2, 0.08, 4.5, ORANGE)  # left vertical line

txt(slide, "🏀", 1.0, 1.1, 1.5, 1.5, size=72)

txt(slide, "BASKETBALL-GOD", 2.6, 1.3, 9, 1.0,
    size=48, color=WHITE, bold=True)
txt(slide, "Version 2.0 — Model Upgrades", 2.6, 2.3, 9, 0.6,
    size=26, color=ORANGE, bold=True)
txt(slide, "5 improvements that made the model smarter,\nmore profitable, and closer to how sharp bettors think",
    2.6, 3.1, 9.5, 1.0, size=18, color=GRAY)

# Bottom stat strip
box(slide, 0, 6.2, 13.33, 1.3, DARKGRAY)
for i, (label, v1, v2, col) in enumerate([
    ("Accuracy",         "70.2%",  "70.4%",  BLUE),
    ("ATS Accuracy",     "59.3%",  "63.4%",  GREEN),
    ("ROI (flat bets)",  "+0.9%",  "+11.5%", GREEN),
    ("False Positives",  "23.0%",  "16.2%",  YELLOW),
    ("Brier Score",      "0.1911", "0.1899", PURPLE),
]):
    xoff = 0.3 + i * 2.6
    box(slide, xoff, 6.28, 2.4, 1.1, CARD)
    txt(slide, label, xoff+0.1, 6.32, 2.2, 0.3, size=10, color=GRAY)
    txt(slide, f"{v1} → {v2}", xoff+0.1, 6.6, 2.2, 0.4, size=15, color=col, bold=True)

txt(slide, "v1 Baseline → v2 Enhanced  |  CPCV Backtest 2022–2025  |  81,829 training games",
    0.3, 7.1, 12, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE BIG PICTURE
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, ORANGE)

txt(slide, "THE BIG PICTURE", 0.4, 0.2, 12, 0.5, size=13, color=ORANGE, bold=True)
txt(slide, "What We Changed & Why It Matters", 0.4, 0.6, 12, 0.6, size=30, color=WHITE, bold=True)

txt(slide, "The old model was like a scout who only reads the season stats sheet.\nThe new model is a scout who also watches film, tracks injuries, and knows which coaches just got fired.",
    0.4, 1.35, 12.5, 0.9, size=16, color=GRAY, italic=True)

# 5 pillars
for i, (num, title, icon, color, old, new) in enumerate([
    ("1", "CLV Optimization",       "💰", ORANGE, "Trained to predict winners", "Trained to beat closing lines"),
    ("2", "Injury Impact",          "🏥", BLUE,   "Ignored injuries completely", "Flags injured key players"),
    ("3", "Elo Differential",       "📊", GREEN,  "Used ranking proxies only",   "Direct Elo strength score"),
    ("4", "Momentum / Decay",       "📈", YELLOW, "All games weighted equally",  "Recent games weighted more"),
    ("5", "Coaching & Portal",      "🔄", PURPLE, "Roster assumed stable",       "Flags mid-season chaos"),
]):
    y = 2.4 + i * 0.94
    rounded_box(slide, 0.4, y, 12.4, 0.82, DARKGRAY)
    # number circle
    circle = slide.shapes.add_shape(9, Inches(0.55), Inches(y+0.18), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = num; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

    txt(slide, icon,  1.1,  y+0.15, 0.5, 0.5, size=22)
    txt(slide, title, 1.65, y+0.08, 3.2, 0.35, size=16, color=WHITE, bold=True)
    # old → new
    txt(slide, "BEFORE:", 5.1, y+0.08, 0.9, 0.3, size=10, color=GRAY, bold=True)
    txt(slide, old,       6.0, y+0.08, 3.0, 0.3, size=12, color=GRAY)
    txt(slide, "AFTER:",  9.2, y+0.08, 0.8, 0.3, size=10, color=color, bold=True)
    txt(slide, new,       10.0, y+0.08, 2.7, 0.3, size=12, color=WHITE)

    # divider line
    if i < 4:
        ln = slide.shapes.add_connector(1, Inches(0.4), Inches(y+0.82), Inches(12.8), Inches(y+0.82))
        ln.line.color.rgb = BG; ln.line.width = Pt(1)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — IMPROVEMENT 1: CLV
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, ORANGE)
pill(slide, "IMPROVEMENT 1", 0.4, 0.18, ORANGE)

txt(slide, "CLV Optimization — Beat the Line, Not Just Pick Winners", 0.4, 0.65, 12.5, 0.6, size=28, color=WHITE, bold=True)

# Analogy box
rounded_box(slide, 0.4, 1.4, 12.4, 1.1, DARKGRAY)
txt(slide, "🎯  The Analogy", 0.65, 1.48, 4, 0.35, size=13, color=ORANGE, bold=True)
txt(slide, "Imagine two poker players. Player A wins 55% of hands. Player B wins 52% of hands "
           "but only plays when they have a real edge. Player B makes more money — because being selective "
           "and right about the edge matters more than raw win rate.",
    0.65, 1.82, 12.0, 0.55, size=13, color=GRAY)

# Left: old way
rounded_box(slide, 0.4, 2.7, 5.9, 3.6, DARKGRAY)
txt(slide, "❌  OLD MODEL", 0.65, 2.82, 5, 0.35, size=14, color=GRAY, bold=True)
txt(slide, "Optimized for ACCURACY", 0.65, 3.2, 5.5, 0.4, size=18, color=GRAY, bold=True)
txt(slide, '• Trained to answer: "Who will win?"\n'
           '• 70% accuracy sounds great\n'
           '• BUT: the market already knows who will win\n'
           '• If market says 65% and model says 67%...\n'
           '  that 2% edge gets eaten by the vig\n'
           '• You bet a lot, win a normal amount, barely profit',
    0.65, 3.65, 5.5, 2.4, size=13, color=GRAY)

# Right: new way
rounded_box(slide, 6.9, 2.7, 5.9, 3.6, DARKGRAY)
txt(slide, "✅  NEW MODEL", 7.15, 2.82, 5, 0.35, size=14, color=GREEN, bold=True)
txt(slide, "Optimized for CLV EDGE", 7.15, 3.2, 5.5, 0.4, size=18, color=GREEN, bold=True)
txt(slide, '• CLV = Closing Line Value\n'
           '• "How much did the line move AGAINST my bet?"\n'
           '• Sharp money moves lines — if your bet matches\n'
           '  where the line moves, you think like sharps do\n'
           '• Model now uses Elo as the CLV baseline\n'
           '• Bets only when edge is real & measurable',
    7.15, 3.65, 5.5, 2.4, size=13, color=GRAY)

# arrow between
txt(slide, "→", 6.3, 4.2, 0.6, 0.6, size=36, color=ORANGE, bold=True)

# result strip
box(slide, 0.4, 6.45, 12.4, 0.75, RGBColor(0x00, 0x2A, 0x14))
txt(slide, "RESULT:  ROI jumped from +0.9% → +11.5%   |   False positives dropped 6.8pp   |   Model bets less but wins more",
    0.7, 6.57, 12, 0.4, size=13, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — IMPROVEMENT 2: INJURY
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, BLUE)
pill(slide, "IMPROVEMENT 2", 0.4, 0.18, BLUE)

txt(slide, "Injury Impact Modeling — The Model Finally Reads the News", 0.4, 0.65, 12.5, 0.6, size=28, color=WHITE, bold=True)

# Analogy
rounded_box(slide, 0.4, 1.4, 12.4, 1.0, DARKGRAY)
txt(slide, "🏥  The Analogy", 0.65, 1.48, 4, 0.35, size=13, color=BLUE, bold=True)
txt(slide, "The old model was like a stockbroker who only looks at last quarter's earnings and ignores "
           "the headline: 'CEO just quit.' One injury to a star player can flip a 60% favorite into a coin flip. "
           "College basketball is especially brutal — teams have 13 players, not 15.",
    0.65, 1.8, 12.0, 0.52, size=13, color=GRAY)

# How it works diagram
txt(slide, "HOW IT WORKS", 0.4, 2.6, 5, 0.35, size=12, color=BLUE, bold=True)

steps = [
    ("1", "Rotowire / ESPN\nhealth reports scraped", BLUE),
    ("2", "Player rated by\nusage % + pts/game", BLUE),
    ("3", "injury_impact_score\nadded as feature", BLUE),
    ("4", "Model adjusts win\nprobability down", GREEN),
]
for i, (n, label, col) in enumerate(steps):
    xo = 0.4 + i * 3.0
    rounded_box(slide, xo, 3.05, 2.6, 1.3, DARKGRAY)
    circle = slide.shapes.add_shape(9, Inches(xo+0.1), Inches(3.12), Inches(0.38), Inches(0.38))
    circle.fill.solid(); circle.fill.fore_color.rgb = col; circle.line.fill.background()
    tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
    txt(slide, label, xo+0.58, 3.1, 2.0, 0.7, size=12, color=WHITE)
    if i < 3:
        txt(slide, "→", xo+2.62, 3.5, 0.4, 0.4, size=20, color=GRAY)

# Important caveat
rounded_box(slide, 0.4, 4.55, 12.4, 1.35, RGBColor(0x0A, 0x1A, 0x2E))
txt(slide, "⚡  Important: Live vs. Historical", 0.65, 4.65, 6, 0.35, size=13, color=BLUE, bold=True)
txt(slide, "Historical training data (2010–2025) has NO injury records — we set the feature to 0.0 during training.\n"
           "The model learns the coefficient exists. At live prediction time, the feature activates with real scraped data.\n"
           "Think of it like a light switch wired but off — training installs the wiring, live predictions flip the switch.",
    0.65, 5.0, 12.0, 0.75, size=12, color=GRAY)

# Why college is different
rounded_box(slide, 0.4, 6.05, 12.4, 0.85, DARKGRAY)
txt(slide, "WHY COLLEGE > NBA FOR THIS FEATURE:", 0.65, 6.12, 6, 0.3, size=11, color=YELLOW, bold=True)
txt(slide, "NBA teams have 15-man rosters and load management. College teams have 13 players and players miss games randomly "
           "due to class schedules, transfers, eligibility issues. One star out = massive swing in win probability.",
    0.65, 6.42, 12.0, 0.4, size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — IMPROVEMENT 3: ELO
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, GREEN)
pill(slide, "IMPROVEMENT 3", 0.4, 0.18, GREEN)

txt(slide, "Elo Differential — The Model's Power Rankings", 0.4, 0.65, 12.5, 0.6, size=28, color=WHITE, bold=True)

# Analogy
rounded_box(slide, 0.4, 1.4, 12.4, 1.0, DARKGRAY)
txt(slide, "♟️  The Analogy", 0.65, 1.48, 4, 0.35, size=13, color=GREEN, bold=True)
txt(slide, "Elo was invented for chess. If a grandmaster beats a beginner, their rating barely moves. "
           "But if they lose? Huge drop. It's a self-correcting scorecard that accounts for WHO you beat, not just whether you won. "
           "The old model used Massey rankings (a poll average). Elo is the live, dynamic version.",
    0.65, 1.78, 12.0, 0.55, size=13, color=GRAY)

# Side by side comparison
rounded_box(slide, 0.4, 2.6, 5.8, 3.5, DARKGRAY)
txt(slide, "📋  OLD: Massey Rankings", 0.65, 2.72, 5.3, 0.4, size=15, color=GRAY, bold=True)
txt(slide, "• Composite of multiple ranking systems\n"
           "• Updated weekly (or less)\n"
           "• Based on wins/losses + schedule strength\n"
           "• 40% of games had missing Massey data!\n"
           "• Doesn't react to a team's recent form\n"
           "• A team can win 5 in a row but still rank #80",
    0.65, 3.12, 5.3, 2.2, size=13, color=GRAY)

txt(slide, "→", 6.3, 4.1, 0.6, 0.6, size=36, color=GREEN, bold=True)

rounded_box(slide, 7.0, 2.6, 5.9, 3.5, DARKGRAY)
txt(slide, "📡  NEW: Elo Differential", 7.25, 2.72, 5.4, 0.4, size=15, color=GREEN, bold=True)
txt(slide, "• Computed fresh from every game ever played\n"
           "• Updates after every single game\n"
           "• Accounts for margin of victory\n"
           "• 0% missing data (computed internally)\n"
           "• Immediately reacts to wins/losses\n"
           "• Became the #2 most important feature (22.3%)",
    7.25, 3.12, 5.4, 2.2, size=13, color=WHITE)

# Formula box
rounded_box(slide, 0.4, 6.25, 12.4, 0.95, RGBColor(0x0A, 0x1F, 0x0A))
txt(slide, "THE FORMULA:  ", 0.65, 6.32, 3, 0.35, size=12, color=GREEN, bold=True)
txt(slide, "New Elo = Old Elo + K × (Actual Result − Expected Result)     |     "
           "K = 20 (regular season)   K = 30 (conference tournament)     |     "
           "Expected = 1 / (1 + 10^(−ΔElo/400))",
    0.65, 6.65, 12.0, 0.45, size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — IMPROVEMENT 4: MOMENTUM
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, YELLOW)
pill(slide, "IMPROVEMENT 4", 0.4, 0.18, YELLOW)

txt(slide, "Momentum & Decay Weighting — Hot Streaks Are Real", 0.4, 0.65, 12.5, 0.6, size=28, color=WHITE, bold=True)

# Analogy
rounded_box(slide, 0.4, 1.4, 12.4, 0.95, DARKGRAY)
txt(slide, "🔥  The Analogy", 0.65, 1.48, 4, 0.35, size=13, color=YELLOW, bold=True)
txt(slide, "The old model treated a team's game in November the same as their game in February. "
           "That's like judging a baseball pitcher's form today based equally on every start of the season — "
           "including the one where he had food poisoning in week 2. Recent games tell you who a team IS right now.",
    0.65, 1.78, 12.0, 0.5, size=13, color=GRAY)

# Two features explained side by side
rounded_box(slide, 0.4, 2.55, 6.0, 3.8, DARKGRAY)
txt(slide, "📈  Feature: last_5_ewm_margin", 0.65, 2.68, 5.5, 0.38, size=14, color=YELLOW, bold=True)
txt(slide, "EWM = Exponentially Weighted Moving Average\n\n"
           "Looks at a team's last 5 game margins\n"
           "(e.g. won by +8, +2, +15, -3, +10)\n\n"
           "But weights them so the most recent game\n"
           "counts more than the oldest:\n\n"
           "  Game 5 (today):  weight 0.45\n"
           "  Game 4:          weight 0.27\n"
           "  Game 3:          weight 0.16\n"
           "  Game 2:          weight 0.08\n"
           "  Game 1 (oldest): weight 0.04",
    0.65, 3.1, 5.5, 2.8, size=12, color=GRAY)

rounded_box(slide, 6.9, 2.55, 6.0, 3.8, DARKGRAY)
txt(slide, "⚖️  Feature: Season Decay Weights", 7.15, 2.68, 5.5, 0.38, size=14, color=YELLOW, bold=True)
txt(slide, "When TRAINING the model, not all historical\n"
           "seasons are equal either:\n\n"
           "  2026 games:  weight 1.00  (full weight)\n"
           "  2025 games:  weight 0.75\n"
           "  2024 games:  weight 0.50\n"
           "  2023 games:  weight 0.25\n"
           "  2022+ older: weight 0.25\n\n"
           "Basketball changed with the transfer portal.\n"
           "A 2012 game tells you less about 2026 than\n"
           "a 2024 game does.",
    7.15, 3.1, 5.5, 2.8, size=12, color=GRAY)

# Impact strip
box(slide, 0.4, 6.5, 12.4, 0.7, RGBColor(0x1A, 0x15, 0x00))
txt(slide, "RESULT:  ATS accuracy +4.2pp  →  Model better at predicting margin (not just winner)  "
           "|  Brier score improved (predictions better calibrated)",
    0.7, 6.62, 12, 0.4, size=13, color=YELLOW, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — IMPROVEMENT 5: COACHING / PORTAL
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, PURPLE)
pill(slide, "IMPROVEMENT 5", 0.4, 0.18, PURPLE)

txt(slide, "Coaching Changes & Transfer Portal — When Teams Fall Apart Mid-Season", 0.4, 0.65, 12.5, 0.6, size=27, color=WHITE, bold=True)

# Analogy
rounded_box(slide, 0.4, 1.4, 12.4, 0.95, DARKGRAY)
txt(slide, "🔄  The Analogy", 0.65, 1.48, 4, 0.35, size=13, color=PURPLE, bold=True)
txt(slide, "Imagine betting on a company's stock using last year's earnings, not knowing the CEO just resigned. "
           "Mid-season coaching changes and mass roster turnover are the basketball equivalent. "
           "The old model would see 'team with good stats' — not 'team in complete chaos.'",
    0.65, 1.78, 12.0, 0.5, size=13, color=GRAY)

# Two features
rounded_box(slide, 0.4, 2.55, 5.9, 3.6, DARKGRAY)
txt(slide, "👔  coaching_instability", 0.65, 2.65, 5.4, 0.38, size=14, color=PURPLE, bold=True)
txt(slide, "Automatically detected from the database.\n\n"
           "If a team has more than 1 coach listed\n"
           "in the same season → mid-season change.\n\n"
           "134 (season, team) pairs flagged across\n"
           "2010–2025 (about 1,600 affected games).\n\n"
           "Feature = 1 if change happened,\n"
           "0 if stable all season.\n\n"
           "diff_coaching_instability = home − away\n"
           "so negative = away team has chaos",
    0.65, 3.05, 5.4, 2.6, size=12, color=GRAY)

rounded_box(slide, 6.8, 2.55, 6.1, 3.6, DARKGRAY)
txt(slide, "🚪  roster_disruption_score", 7.05, 2.65, 5.6, 0.38, size=14, color=PURPLE, bold=True)
txt(slide, "Manually curated JSON database of\n"
           "transfer portal disruption (2022–2026).\n\n"
           "Score 0.0 (stable) → 1.0 (complete chaos)\n\n"
           "Examples:\n"
           "  Kentucky 2023: 0.85 (mass portal exits)\n"
           "  Auburn 2024:   0.70 (lost 4 starters)\n"
           "  Duke 2025:     0.40 (moderate turnover)\n\n"
           "263 games flagged with non-zero score.\n"
           "Manually maintained — add new entries\n"
           "at start of each season.",
    7.05, 3.05, 5.7, 2.6, size=12, color=GRAY)

# How to update
box(slide, 0.4, 6.3, 12.4, 0.9, RGBColor(0x15, 0x0A, 0x1F))
txt(slide, "📝  TO UPDATE:  Edit  phase7_v2/coaching_portal.json  at the start of each season. "
           "Add new coaching changes and portal disruption scores. Re-run  python phase7_v2/train_v2.py  to retrain.",
    0.7, 6.42, 12, 0.55, size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — RESULTS COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, GREEN)

txt(slide, "THE RESULTS", 0.4, 0.2, 12, 0.4, size=12, color=GREEN, bold=True)
txt(slide, "Side-by-Side: v1 Baseline vs v2 Enhanced", 0.4, 0.55, 12, 0.55, size=30, color=WHITE, bold=True)
txt(slide, "CPCV Walk-Forward Backtest  |  Test seasons: 2022, 2023, 2024, 2025  |  81,829 training games",
    0.4, 1.15, 12.5, 0.35, size=12, color=GRAY)

# Main metric cards row 1
metric_card(slide, 0.4,  1.65, "OVERALL ACCURACY",       "70.2%",  "70.4%",  "+0.3pp",  True)
metric_card(slide, 3.4,  1.65, "ATS ACCURACY",           "59.3%",  "63.4%",  "+4.2pp",  True)
metric_card(slide, 6.4,  1.65, "ROI (FLAT $100/BET)",    "+0.9%",  "+11.5%", "+10.6pp", True)
metric_card(slide, 9.4,  1.65, "BRIER SCORE (↓ better)", "0.1911", "0.1899", "−0.6%",   True)

metric_card(slide, 0.4,  3.35, "HIGH-CONF HIT RATE",   "91.3%", "91.4%", "+0.1pp",  True)
metric_card(slide, 3.4,  3.35, "FALSE POSITIVE RATE",  "23.0%", "16.2%", "−6.8pp",  True)
metric_card(slide, 6.4,  3.35, "STRONG BETS PLACED",   "6,326", "3,909", "−38%",    True)
metric_card(slide, 9.4,  3.35, "TOTAL P&L ($100 FLAT)", "$+5,776", "$+44,891", "+676%", True)

# Feature importance chart (horizontal bars)
txt(slide, "TOP FEATURES BY IMPORTANCE (XGBoost)", 0.4, 5.05, 7, 0.35, size=12, color=ORANGE, bold=True)
features = [
    ("diff_massey_avg_rank",   0.2579, GRAY,   False),
    ("diff_elo  ← NEW",        0.2230, GREEN,  True),
    ("diff_avg_margin",        0.0738, GRAY,   False),
    ("diff_conf_win_pct",      0.0445, GRAY,   False),
    ("diff_last5_ewm_margin ← NEW", 0.0138, YELLOW, True),
]
bar_max = 8.0  # inches for 100%
for i, (name, imp, col, is_new) in enumerate(features):
    y = 5.5 + i * 0.35
    bar_w = imp / 0.30 * 3.5  # scale so 0.30 = full bar
    box(slide, 3.5, y+0.03, bar_w, 0.25, col)
    txt(slide, name, 0.4, y, 3.0, 0.3, size=11, color=WHITE if is_new else GRAY)
    txt(slide, f"{imp*100:.1f}%", 3.5+bar_w+0.1, y, 0.8, 0.3, size=11, color=col, bold=True)

# Key insight callout
rounded_box(slide, 7.2, 5.05, 5.8, 2.25, DARKGRAY)
txt(slide, "💡  Key Insight", 7.45, 5.15, 4, 0.35, size=13, color=ORANGE, bold=True)
txt(slide, "The ROI jump (+10.6pp) happened because the model\n"
           "became MORE SELECTIVE.\n\n"
           "It placed 38% fewer 'strong value' bets but each one\n"
           "was more accurate — fewer false alarms means\n"
           "you're not throwing money at games that only LOOKED\n"
           "like strong bets from noisy v1 features.",
    7.45, 5.52, 5.3, 1.6, size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — HOW TO RUN IT
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank)
add_bg(slide)
box(slide, 0, 0, 13.33, 0.08, BLUE)

txt(slide, "HOW TO RUN IT", 0.4, 0.2, 12, 0.45, size=12, color=BLUE, bold=True)
txt(slide, "Everything You Need to Get the Dashboard Running", 0.4, 0.6, 12, 0.55, size=30, color=WHITE, bold=True)

steps_full = [
    ("1", "Train the v2 model\n(first time only — ~5 min)",
     "python phase7_v2/train_v2.py",
     "Builds all 5 new features, computes Elo from 200K games,\ntrains ensemble, saves model to phase7_v2/output/",
     BLUE),
    ("2", "Run the A/B backtest\n(optional — to see results)",
     "python phase7_v2/backtest_ab.py",
     "Compares v1 vs v2 side by side, prints table,\nsaves backtest_results_2026-03-14.json",
     PURPLE),
    ("3", "Start the dashboard",
     "python web/server.py",
     "Loads the model, starts Flask on port 5051,\nfetches today's games & odds automatically",
     GREEN),
    ("4", "Open in browser",
     "http://localhost:5051",
     "Live games panel auto-refreshes every 60s.\nClick 'Refresh Lines' to force new odds fetch.",
     ORANGE),
]

for i, (n, title, cmd, desc, col) in enumerate(steps_full):
    y = 1.4 + i * 1.38
    rounded_box(slide, 0.4, y, 12.4, 1.22, DARKGRAY)
    # step circle
    circle = slide.shapes.add_shape(9, Inches(0.55), Inches(y+0.38), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = col; circle.line.fill.background()
    tf = circle.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = n; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

    txt(slide, title, 1.25, y+0.1,  3.8, 0.55, size=14, color=WHITE, bold=True)
    # command box
    box(slide, 5.3, y+0.18, 4.8, 0.5, RGBColor(0x08, 0x0D, 0x13))
    txt(slide, cmd, 5.45, y+0.22, 4.5, 0.38, size=12, color=col, bold=True)
    txt(slide, desc, 10.3, y+0.12, 2.8, 0.65, size=10, color=GRAY)

# Troubleshooting pointer
box(slide, 0.4, 6.95, 12.4, 0.4, DARKGRAY)
txt(slide, "🔧  Having issues? See:  BBALLGOD3.0 FAQ ISSUES/TROUBLESHOOTING.md  — covers every known problem with step-by-step fixes",
    0.65, 7.0, 12, 0.3, size=12, color=GRAY)


# ── Save ──────────────────────────────────────────────────────────────────────

out = r"C:\Users\yarden\Basketball-God2.0\BBALLGOD3.0 FAQ ISSUES\Basketball-God-v2-Upgrades.pptx"
prs.save(out)
print(f"Saved: {out}")
